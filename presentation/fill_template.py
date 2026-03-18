"""
Fill the GEF C3 Pitch Deck template with Dream Weaver content.

Run:  bedtime/bin/python presentation/fill_template.py
Output: presentation/DreamWeaver_Final.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "template.pptx")
OUTPUT   = os.path.join(HERE, "DreamWeaver_Final.pptx")

PHOTO_TAMAS     = os.path.join(HERE, "photo_tamas.jpeg")
PHOTO_ALESSANDRO = os.path.join(HERE, "photo_alessandro.jpeg")


def set_text(shape, text, font_size=None, bold=None, color=None):
    """Replace ALL text in a shape, preserving first-run formatting."""
    tf = shape.text_frame
    # Clear all paragraphs except the first
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    # Clear runs
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    p.text = text
    if font_size is not None:
        p.font.size = Pt(font_size)
    if bold is not None:
        p.font.bold = bold
    if color is not None:
        p.font.color.rgb = color


def set_multiline(shape, lines, font_size=12, color=None, bold=False):
    """Replace shape text with multiple lines."""
    tf = shape.text_frame
    # Clear existing
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        if color:
            p.font.color.rgb = color
        p.font.bold = bold
        p.space_before = Pt(4)


def find_shape(slide, name):
    """Find a shape by its name."""
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def find_shape_by_text(slide, text):
    """Find a shape containing specific text."""
    for s in slide.shapes:
        if s.has_text_frame and text.lower() in s.text_frame.text.lower():
            return s
    return None


# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation(TEMPLATE)

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
sl = prs.slides[1]
for s in sl.shapes:
    if s.has_text_frame:
        if "GROUP" in s.text_frame.text.upper():
            set_text(s, "Team 2 — GEF C3")
        elif "TITLE" in s.text_frame.text.lower() or "project" in s.text_frame.text.lower():
            set_text(s, "Dream Weaver")


# ── SLIDE 2: Team ─────────────────────────────────────────────────────────────
sl = prs.slides[2]

# Update the title
title_shape = find_shape_by_text(sl, "Meet The Team")
if title_shape:
    set_text(title_shape, "Meet The Team ✨")

# Team data: (name, role, job_profile, photo_path_or_None)
team = [
    ("Tamas", "Backend Architect", "Backend Developer — DERTOUR", PHOTO_TAMAS),
    ("Alessandro", "PM / Team Lead", "AI Engineer — Rebis Labs", PHOTO_ALESSANDRO),
    ("Om", "Frontend Dev", "Configuration Page & UI", None),
    ("Kumaraguru", "Prompt Engineer", "Moral Lessons & Prompts", None),
    ("Ravi", "Backend / RAG", "FAISS Research & Backend", None),
]

# The template has 4 member slots. Each slot consists of:
#   - A photo (PICTURE shape)
#   - A name text box
#   - A "brief" text box
#   - A "job profile" text box
#   - A circle shape (decorative)

# Collect all relevant shapes by their approximate X position
# Photos are at x positions: ~1181526, ~3065915, ~4962537, ~6846927
# That's roughly 4 columns spaced ~1.87M apart

# Find existing photo shapes (PICTUREs with specific sizes ~1115700x1064700)
photo_shapes = []
name_shapes = []
brief_shapes = []
job_shapes = []

for s in sl.shapes:
    if not s.has_text_frame:
        # Check if it's a team photo placeholder (specific size)
        if abs(s.width - 1115700) < 50000 and abs(s.height - 1064700) < 50000:
            photo_shapes.append(s)
    else:
        txt = s.text_frame.text.strip()
        if txt in ("Lucy", "Dan", "Alex", "Stephy"):
            name_shapes.append(s)
        elif "Brief about" in txt:
            brief_shapes.append(s)
        elif "job profile" in txt:
            job_shapes.append(s)

# Sort all by left position
photo_shapes.sort(key=lambda s: s.left)
name_shapes.sort(key=lambda s: s.left)
brief_shapes.sort(key=lambda s: s.left)
job_shapes.sort(key=lambda s: s.left)

print(f"Found: {len(photo_shapes)} photos, {len(name_shapes)} names, "
      f"{len(brief_shapes)} briefs, {len(job_shapes)} jobs")

# Fill first 4 slots
for i in range(min(4, len(team))):
    name, role, job_profile, photo = team[i]

    if i < len(name_shapes):
        set_text(name_shapes[i], name, font_size=14, bold=True)
    if i < len(brief_shapes):
        set_text(brief_shapes[i], role, font_size=10)
    if i < len(job_shapes):
        set_text(job_shapes[i], job_profile, font_size=9)

    # Replace photo if available
    if photo and os.path.exists(photo) and i < len(photo_shapes):
        old_photo = photo_shapes[i]
        # Add new image at same position/size
        pic = sl.shapes.add_picture(
            photo,
            old_photo.left, old_photo.top,
            old_photo.width, old_photo.height,
        )
        # Remove old placeholder image
        sp = old_photo._element
        sp.getparent().remove(sp)

# Add 5th member — squeeze in by adjusting positions
# We need to make room. Let's shift everyone to be evenly spaced for 5.
# Template slide width = 9144000 EMU (~10 inches)
# Photo width = 1115700 EMU (~1.22 inches)
# For 5 photos: total_spacing = 9144000 - 5*1115700 = 3565500
# margin = 3565500 / 6 = 594250

if len(team) > 4:
    # Recalculate positions for all 5
    photo_w = 1115700
    n_members = 5
    total_space = 9144000 - n_members * photo_w
    gap = total_space // (n_members + 1)

    # Gather ALL shapes per column (by original x proximity)
    # First, rebuild lists of current shapes
    all_shapes_by_col = []  # list of (photo, name, brief, job) per column

    # Re-find shapes after modifications
    current_photos = []
    current_names = []
    current_briefs = []
    current_jobs = []

    for s in sl.shapes:
        if not s.has_text_frame:
            if abs(s.width - 1115700) < 50000 and abs(s.height - 1064700) < 50000:
                current_photos.append(s)
            # Also check for newly added pictures
            elif hasattr(s, 'image') and s.width > 500000:
                current_photos.append(s)
        else:
            txt = s.text_frame.text.strip()
            if txt in [t[0] for t in team[:4]]:
                current_names.append(s)
            elif txt in [t[1] for t in team[:4]]:
                current_briefs.append(s)
            elif txt in [t[2] for t in team[:4]]:
                current_jobs.append(s)

    current_photos.sort(key=lambda s: s.left)
    current_names.sort(key=lambda s: s.left)
    current_briefs.sort(key=lambda s: s.left)
    current_jobs.sort(key=lambda s: s.left)

    # Reposition existing 4 columns
    for i in range(min(4, len(current_photos))):
        new_x = gap + i * (photo_w + gap)
        dx = new_x - current_photos[i].left
        current_photos[i].left = new_x
        if i < len(current_names):
            current_names[i].left += dx
        if i < len(current_briefs):
            current_briefs[i].left += dx
        if i < len(current_jobs):
            current_jobs[i].left += dx

    # Also reposition any decorative circles
    for s in sl.shapes:
        if s.shape_type == 1 and not s.has_text_frame:  # AUTO_SHAPE
            # Find which column this belongs to by proximity
            pass  # Skip decorative repositioning for simplicity

    # Add 5th member
    fifth = team[4]
    x5 = gap + 4 * (photo_w + gap)

    # Use reference positions from first column for Y coordinates
    photo_y = 1406100
    name_y = photo_y + 1064700 + Pt(8)  # below photo
    brief_y = name_y + Pt(20)
    job_y = brief_y + Pt(16)

    # Add name
    from pptx.util import Emu
    name_box = sl.shapes.add_textbox(x5, name_y, photo_w, Pt(20))
    p = name_box.text_frame.paragraphs[0]
    p.text = fifth[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add role
    role_box = sl.shapes.add_textbox(x5, brief_y, photo_w, Pt(16))
    p = role_box.text_frame.paragraphs[0]
    p.text = fifth[1]
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # Add job
    job_box = sl.shapes.add_textbox(x5, job_y, photo_w, Pt(14))
    p = job_box.text_frame.paragraphs[0]
    p.text = fifth[2]
    p.font.size = Pt(9)
    p.alignment = PP_ALIGN.CENTER

    # Add placeholder circle for missing photo
    from pptx.enum.shapes import MSO_SHAPE
    circle = sl.shapes.add_shape(
        MSO_SHAPE.OVAL, x5, photo_y, photo_w, 1064700
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    circle.line.fill.background()
    # Add initials
    ctf = circle.text_frame
    ctf.paragraphs[0].text = "R"
    ctf.paragraphs[0].font.size = Pt(36)
    ctf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    from pptx.enum.text import MSO_ANCHOR
    ctf.word_wrap = False


# ── SLIDE 3: Problem ──────────────────────────────────────────────────────────
sl = prs.slides[3]
for s in sl.shapes:
    if not s.has_text_frame:
        continue
    txt = s.text_frame.text.strip()

    if "What is the Problem" in txt:
        set_text(s, "What is the Problem?", font_size=24, bold=True)
    elif "one line Problem" in txt:
        set_text(s, "Parents are exhausted at bedtime — and no tool creates personalized, illustrated, narrated, interactive stories for their child in real time.", font_size=14)
    elif "Supporting Point - 1" in txt:
        set_multiline(s, [
            "No Personalization",
            "",
            "Generic storybooks don't feature the child's name, pet, or family members",
        ], font_size=11)
    elif "Supporting Point - 2" in txt:
        set_multiline(s, [
            "Repetitive Content",
            "",
            "Reading the same 5 books every night — kids lose interest, parents run out of ideas",
        ], font_size=11)
    elif "Supporting Point - 3" in txt:
        set_multiline(s, [
            "AI Tools Fall Short",
            "",
            "Existing AI generates text-only walls — no illustrations, no voice, no interactivity",
        ], font_size=11)


# ── SLIDE 4: Solution ─────────────────────────────────────────────────────────
sl = prs.slides[4]
for s in sl.shapes:
    if not s.has_text_frame:
        continue
    txt = s.text_frame.text.strip()

    if "How are you solving" in txt:
        set_text(s, "How are we solving it?", font_size=24, bold=True)
    elif "one line solution" in txt:
        set_text(s, "Dream Weaver generates personalized bedtime stories where the child is the hero — with AI illustrations, voice narration, and interactive choices.", font_size=14)
    elif "Must have" in txt:
        set_multiline(s, [
            "Must Have ✓",
            "",
            "• Personalized story (child as hero)",
            "• AI illustrations per scene",
            "• Expressive voice narration",
            "• Interactive choices (2 per scene)",
            "• Safety content filter",
            "• Moral lesson in every story",
        ], font_size=10)
    elif "Should have" in txt:
        set_multiline(s, [
            "Should Have ✓",
            "",
            "• Side characters (family, pets)",
            "• Character visual consistency",
            "• PDF booklet export",
            "• RAG story memory",
        ], font_size=10)
    elif "Could have" in txt:
        set_multiline(s, [
            "Could Have",
            "",
            "• Story library (browse past stories)",
            "• Multi-language narration",
            "• Go-back and re-choose",
            "• Online deployment",
        ], font_size=10)


# ── SLIDE 5: Target Audience ──────────────────────────────────────────────────
sl = prs.slides[5]
for s in sl.shapes:
    if not s.has_text_frame:
        continue
    txt = s.text_frame.text.strip()

    if txt == "Target Audience":
        set_text(s, "Target Audience", font_size=28, bold=True)
    elif "Primary" in txt:
        set_multiline(s, [
            "Primary: Parents of children 3-8",
            "",
            "• Tired parents who want unique, engaging bedtime stories",
            "• Want their child to be the hero of the story",
            "• Value personalization: child's name, pet, family woven in",
            "• Appreciate moral lessons (kindness, sharing, empathy)",
            "",
            "Key pain: exhausted by 8 PM, need something fresh every night",
        ], font_size=11)
    elif "Secondary" in txt:
        set_multiline(s, [
            "Secondary: Educators & Caregivers",
            "",
            "• Teachers looking for personalized reading material",
            "• Childcare providers needing engaging story time",
            "• Grandparents / extended family during visits",
            "",
            "Future: Children's therapists (therapeutic storytelling)",
        ], font_size=11)


# ── SLIDE 6: Product Link ────────────────────────────────────────────────────
sl = prs.slides[6]
for s in sl.shapes:
    if s.has_text_frame and "Product Link" in s.text_frame.text:
        set_multiline(s, [
            "Product Links",
            "",
            "GitHub: github.com/Outskill-CH3-Team-2/bedtime-stories-capstone",
            "",
            "Tech Stack: React + FastAPI + LangGraph + FAISS + OpenRouter",
            "Models: GPT-4o (text) • Gemini Flash (images) • GPT-4o Audio (TTS)",
            "",
            "34/34 tests passing • 5 AI models • 8-scene stories • ~20s per scene",
        ], font_size=14)


# ── SLIDE 7: Thank You ───────────────────────────────────────────────────────
# Keep as is, it already says "Thank You ✨"


# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✓ Saved: {OUTPUT}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Upload to Google Drive → auto-converts to Slides")
